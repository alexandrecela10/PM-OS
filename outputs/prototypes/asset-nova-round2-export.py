#!/usr/bin/env python3
"""Build the Asset Nova round-2 deck and export its PPTX."""

from html import escape
from pathlib import Path
import re
import shutil
import subprocess
import sys


ROOT = Path(__file__).resolve().parents[2]
SOURCE = ROOT / "outputs/asset-nova-round2-deck.md"
HTML = ROOT / "outputs/prototypes/asset-nova-round2-deck.html"
PPTX = ROOT / "outputs/prototypes/asset-nova-round2-deck.pptx"


CSS = r"""  :root{
    --bg:#0a0e14; --panel:#11161f; --panel2:#161d29; --line:#1f2937;
    --txt:#e5eaf2; --dim:#8b95a7; --faint:#5b6474;
    --accent:#4f8ef7; --green:#34d399; --amber:#fbbf24; --red:#f87171;
    --purple:#a78bfa; --cyan:#22d3ee;
    --mono:'SF Mono',ui-monospace,Menlo,monospace;
    --sans:-apple-system,'Inter',system-ui,sans-serif;
  }
  *{box-sizing:border-box;margin:0;padding:0}
  html,body{height:100%}
  body{background:var(--bg);color:var(--txt);font-family:var(--sans);overflow:hidden}
  .slide{display:none;position:absolute;inset:0;padding:56px 84px 72px;flex-direction:column;overflow:hidden}
  .slide.active{display:flex;animation:fade .45s}
  @keyframes fade{from{opacity:0;transform:translateY(10px)}to{opacity:1;transform:none}}
  .kicker{font-family:var(--mono);font-size:12px;letter-spacing:2px;text-transform:uppercase;color:var(--accent);margin-bottom:14px}
  h1{font-size:52px;line-height:1.1;font-weight:700;letter-spacing:-.5px}
  h2{font-size:34px;line-height:1.2;font-weight:700;letter-spacing:-.3px;margin-bottom:24px}
  h2 em{font-style:normal;color:var(--accent)}
  .sub{font-size:19px;color:var(--dim);margin-top:16px;line-height:1.5;max-width:900px}
  .body{flex:1;display:flex;flex-direction:column;justify-content:center}
  table{border-collapse:collapse;width:100%;font-size:15px}
  th{font-size:11px;text-transform:uppercase;letter-spacing:1px;color:var(--faint);text-align:left;padding:8px 14px;border-bottom:1px solid var(--line)}
  td{padding:11px 14px;border-bottom:1px solid var(--line);color:var(--dim);line-height:1.45;vertical-align:top}
  td b,td strong{color:var(--txt)}
  tr.hl td{background:#101a2c;color:var(--txt)}
  .cols{display:flex;gap:28px}
  .col{flex:1}
  .card{background:var(--panel);border:1px solid var(--line);border-radius:12px;padding:20px 24px}
  .card h3{font-size:13px;text-transform:uppercase;letter-spacing:1px;color:var(--accent);margin-bottom:10px}
  .card p,.card li{font-size:14.5px;color:var(--dim);line-height:1.55}
  .card ul{padding-left:18px}
  .bignum{display:flex;gap:44px;margin:26px 0}
  .bignum div b{display:block;font-size:46px;font-family:var(--mono);color:var(--amber)}
  .bignum div i{font-style:normal;font-size:13.5px;color:var(--dim);display:block;max-width:220px;margin-top:6px;line-height:1.4}
  .quote{border-left:3px solid var(--accent);padding:14px 22px;background:var(--panel);border-radius:0 10px 10px 0;font-size:16.5px;color:var(--dim);line-height:1.6;max-width:980px}
  .quote b{color:var(--txt)}
  .quote .attr{display:block;margin-top:10px;font-size:13px;color:var(--faint)}
  pre.diagram{font-family:var(--mono);font-size:13.5px;line-height:1.45;color:var(--dim);background:var(--panel);border:1px solid var(--line);border-radius:12px;padding:22px 26px;overflow:auto}
  pre.diagram b{color:var(--txt)}
  .pill{display:inline-block;font-size:11px;padding:3px 10px;border-radius:12px;background:#1a2333;color:var(--accent);border:1px solid #24344f;margin-right:8px}
  .foot{position:absolute;bottom:22px;left:84px;right:84px;display:flex;justify-content:space-between;font-family:var(--mono);font-size:11px;color:var(--faint)}
  .nav-hint{position:fixed;bottom:22px;right:26px;font-family:var(--mono);font-size:11px;color:var(--faint)}
  .progress{position:fixed;top:0;left:0;height:2px;background:var(--accent);transition:width .3s}
  .fromto td:first-child{color:var(--red)}
  .fromto td:last-child{color:var(--green)}
  .green{color:var(--green)} .amber{color:var(--amber)} .red{color:var(--red)} .purple{color:var(--purple)} .cyan{color:var(--cyan)}
  .demo-link{display:inline-block;margin-top:26px;background:var(--accent);color:#fff;font-weight:600;padding:13px 26px;border-radius:10px;text-decoration:none;font-size:16px}
  .steps{display:flex;gap:16px;margin-top:30px}
  .steps .s{flex:1;background:var(--panel);border:1px solid var(--line);border-radius:12px;padding:16px 18px}
  .steps .s b{display:block;font-family:var(--mono);color:var(--accent);font-size:12px;margin-bottom:8px}
  .steps .s p{font-size:13.5px;color:var(--dim);line-height:1.5}
  .appendix-tag{position:absolute;top:22px;right:84px;font-family:var(--mono);font-size:11px;color:var(--purple);letter-spacing:2px}
  .markdown-list{padding-left:24px;color:var(--dim);font-size:17px;line-height:1.5}
  .markdown-list li{margin-bottom:8px}
  .slide-1 .body{justify-content:center}
  .slide-1 h1{font-size:58px}
  .slide-1 .markdown-p:first-of-type{font-size:24px;color:var(--txt);margin-top:20px}
  .slide-2 .markdown-list{font-size:17px;line-height:1.42}
  .slide-4 .markdown-list{font-size:18px;line-height:1.48}
  .slide-5 pre.diagram,.slide-6 pre.diagram{font-size:12px;line-height:1.3;padding:18px 22px;white-space:pre;overflow:hidden}
  .slide-8 table,.slide-10 table,.slide-12 table{font-size:14px}
  .slide-8 .roi-table{table-layout:fixed;font-size:12px}
  .slide-8 .roi-table th,.slide-8 .roi-table td{padding:5px 8px}
  .slide-8 .roi-table th:first-child,.slide-8 .roi-table td:first-child{width:48%}
  .slide-8 .roi-table th:nth-child(2),.slide-8 .roi-table td:nth-child(2){width:24%}
  .slide-8 .roi-table th:nth-child(3),.slide-8 .roi-table td:nth-child(3){width:28%}
  .roi-table tr.money td:not(:first-child){font-family:var(--mono);color:var(--amber);white-space:nowrap}
  .slide-8 .roi-table tr.money td:not(:first-child){white-space:normal}
  .roi-table tr.money td b,.roi-table tr.money td strong{color:var(--amber)}
  .roi-table tr.verdict td{background:#13261f;color:var(--green)}
  .roi-table tr.verdict td:nth-child(4){background:#2b2410;color:var(--amber)}
  .roi-table tr.verdict td b,.roi-table tr.verdict td strong{color:inherit}
  .slide-8 .markdown-p,.slide-10 .markdown-p,.slide-12 .markdown-p{font-size:15px;line-height:1.4}
  .slide-8 .markdown-p:not(:first-child),.slide-10 .markdown-p:not(:first-child){margin-top:10px}
  .slide-14 .quote{margin-bottom:10px;max-width:1180px;font-size:14px;line-height:1.38;padding:10px 18px}
  .slide-15 pre.diagram{font-size:12.5px;line-height:1.34;padding:18px 22px}
  .slide-16 .markdown-p,.slide-17 .markdown-p,.slide-18 .markdown-list,.slide-19 .markdown-list{font-size:16px;line-height:1.48}
  .markdown-p{font-size:17px;color:var(--dim);line-height:1.5;margin-top:12px;max-width:1180px}
  .markdown-p strong{color:var(--txt)}
  .markdown-p em{color:var(--dim)}
  @media print{
    @page{size:338.7mm 190.5mm;margin:0}
    html,body{overflow:visible;height:auto}
    .slide{display:flex !important;position:relative;inset:auto;width:338.7mm;height:190.5mm;page-break-after:always;animation:none}
    .nav-hint,.progress{display:none}
    .foot{position:absolute}
  }"""


JS = r"""const slides=[...document.querySelectorAll('.slide')];
let cur=0;
function show(n){
  cur=Math.max(0,Math.min(slides.length-1,n));
  slides.forEach((s,i)=>s.classList.toggle('active',i===cur));
  document.getElementById('prog').style.width=((cur+1)/slides.length*100)+'%';
}
document.addEventListener('keydown',e=>{
  if(e.key==='ArrowRight'||e.key===' '||e.key==='PageDown')show(cur+1);
  else if(e.key==='ArrowLeft'||e.key==='PageUp')show(cur-1);
  else if(e.key==='Home')show(0);
  else if(e.key==='f'||e.key==='F')document.documentElement.requestFullscreen?.();
});
document.addEventListener('click',e=>{if(e.target.closest('a'))return;show(cur+1)});
const m=location.hash.match(/slide=(\d+)/);show(m?+m[1]:0);
if(location.hash.includes('export'))document.querySelector('.nav-hint').style.display='none';"""


def inline(text: str) -> str:
    value = escape(text, quote=False)
    value = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", value)
    value = re.sub(r"\*(.+?)\*", r"<em>\1</em>", value)
    value = re.sub(r"`(.+?)`", r"<code>\1</code>", value)
    return value


def table_html(lines: list[str], slide_num: int) -> str:
    rows = []
    for line in lines:
        cells = [cell.strip() for cell in line.strip().strip("|").split("|")]
        if cells and not all(set(cell) <= {"-", ":"} for cell in cells):
            rows.append(cells)
    if not rows:
        return ""
    classes = ["roi-table"] if slide_num in {8, 10, 12} else []
    out = [f'<table class="{" ".join(classes)}">']
    headers, data = rows[0], rows[1:]
    out.append("<tr>" + "".join(f"<th>{inline(cell)}</th>" for cell in headers) + "</tr>")
    for cells in data:
        first = cells[0] if cells else ""
        row_classes = []
        if first.lower() in {"total per project", "return", "verdict"}:
            row_classes.append("hl")
        if slide_num in {8, 10, 12} and any(
            word in first.lower() for word in ("claims submitted", "settled", "defence", "damages", "dispute", "price", "rework", "share caused", "catch half", "value per")
        ):
            row_classes.append("money")
        if first.lower() == "verdict":
            row_classes.append("verdict")
        attrs = f' class="{" ".join(row_classes)}"' if row_classes else ""
        out.append(f"<tr{attrs}>" + "".join(f"<td>{inline(cell)}</td>" for cell in cells) + "</tr>")
    out.append("</table>")
    return "\n".join(out)


def render_blocks(body: str, slide_num: int) -> str:
    lines = body.strip().splitlines()
    out = []
    i = 0
    while i < len(lines):
        line = lines[i]
        if not line.strip():
            i += 1
            continue
        if line.strip() == "```":
            i += 1
            code = []
            while i < len(lines) and lines[i].strip() != "```":
                code.append(lines[i])
                i += 1
            i += 1
            code_markup = escape("\n".join(code))
            out.append(f'<pre class="diagram">{code_markup}</pre>')
            continue
        if line.startswith("|"):
            table = []
            while i < len(lines) and lines[i].startswith("|"):
                table.append(lines[i])
                i += 1
            out.append(table_html(table, slide_num))
            continue
        if line.startswith(">"):
            quotes = []
            while i < len(lines) and lines[i].startswith(">"):
                quotes.append(lines[i][1:].strip())
                i += 1
            if slide_num == 14 and len(quotes) >= 2:
                quote_text = inline(quotes[0])
                attr = inline(" ".join(quotes[1:]))
                out.append(f'<div class="quote"><span>{quote_text}</span><span class="attr">{attr}</span></div>')
            else:
                out.append(f'<div class="quote">{inline(" ".join(quotes))}</div>')
            continue
        if re.match(r"^\d+\.\s+", line):
            items = []
            while i < len(lines) and re.match(r"^\d+\.\s+", lines[i]):
                items.append(re.sub(r"^\d+\.\s+", "", lines[i]))
                i += 1
            out.append('<ol class="markdown-list">' + "".join(f"<li>{inline(item)}</li>" for item in items) + "</ol>")
            continue
        if line.startswith("- "):
            items = []
            while i < len(lines) and lines[i].startswith("- "):
                items.append(lines[i][2:])
                i += 1
            out.append('<ul class="markdown-list">' + "".join(f"<li>{inline(item)}</li>" for item in items) + "</ul>")
            continue
        paragraph = [line]
        i += 1
        while i < len(lines) and lines[i].strip() and not (
            lines[i].startswith("|") or lines[i].startswith(">") or lines[i].startswith("- ")
            or re.match(r"^\d+\.\s+", lines[i]) or lines[i].strip() == "```"
        ):
            paragraph.append(lines[i])
            i += 1
        out.append(f'<p class="markdown-p">{inline(" ".join(paragraph))}</p>')
    return "\n".join(out)


def build_html() -> str:
    source = SOURCE.read_text()
    sections = re.split(r"(?m)^## SLIDE (\d+) — (.+)$", source)
    slides = []
    for position in range(1, len(sections), 3):
        slide_num = int(sections[position])
        title = sections[position + 1].strip()
        body = sections[position + 2].split("\n---", 1)[0].strip()
        if slide_num == 1:
            title = "Asset Nova"
            body = re.sub(r"^\*\*Asset Nova\*\*\n?", "", body)
        content = render_blocks(body, slide_num)
        title_markup = f"<h1>{inline(title)}</h1>" if slide_num == 1 else f"<h2>{inline(title)}</h2>"
        slides.append(
            f'<section class="slide slide-{slide_num}{" active" if slide_num == 1 else ""}">'
            f'<div class="kicker">Asset Nova · Round 2 · Slide {slide_num}</div>'
            f'{title_markup}<div class="body">{content}</div>'
            f'<div class="foot"><span>Asset Nova · Round 2</span><span>{slide_num} / 19</span></div></section>'
        )
    return (
        '<!DOCTYPE html>\n<html lang="en">\n<head>\n<meta charset="UTF-8">\n'
        '<meta name="viewport" content="width=device-width, initial-scale=1.0">\n'
        '<title>Asset Nova — Round 2 Deck</title>\n<style>\n' + CSS +
        '\n</style>\n</head>\n<body>\n<div class="progress" id="prog"></div>\n' +
        "\n".join(slides) +
        '\n<div class="nav-hint">← → navigate · F fullscreen</div>\n<script>\n' + JS +
        '\n</script>\n</body>\n</html>\n'
    )


def export_pptx() -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image
        from playwright.sync_api import sync_playwright
    except ImportError as exc:
        raise SystemExit(f"Missing export dependency: {exc}")

    screenshots = ROOT / "outputs/prototypes/asset-nova-round2-screenshots"
    screenshots.mkdir(exist_ok=True)
    url = HTML.resolve().as_uri()
    with sync_playwright() as playwright:
        chrome = shutil.which("google-chrome") or shutil.which("chromium") or shutil.which("chromium-browser")
        browser = playwright.chromium.launch(headless=True, executable_path=chrome)
        page = browser.new_page(viewport={"width": 1440, "height": 900}, device_scale_factor=1)
        for slide_num in range(1, 20):
            page.goto(f"{url}#slide={slide_num - 1}&export", wait_until="load")
            page.screenshot(path=str(screenshots / f"slide-{slide_num:02d}.png"), full_page=False)
        browser.close()

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for slide_num in range(1, 20):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(screenshots / f"slide-{slide_num:02d}.png"),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
    prs.save(PPTX)


if __name__ == "__main__":
    HTML.write_text(build_html())
    if "--html-only" not in sys.argv:
        export_pptx()
    print(HTML)
    if "--html-only" not in sys.argv:
        print(PPTX)
